import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 16:9 Widescreen dimensions
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Premium Modern Theme Palette (Obsidian & Indigo/Violet + Emerald accents)
COLOR_BG = RGBColor(10, 15, 30)          # Deep Obsidian Navy #0a0f1e
COLOR_SURFACE = RGBColor(18, 25, 45)     # Sleek Slate #12192d
COLOR_CARD_BORDER = RGBColor(40, 53, 86) # Subtle Slate Border #283556
COLOR_VIOLET = RGBColor(124, 58, 237)    # Electric Violet #7c3aed
COLOR_CYAN = RGBColor(14, 165, 233)      # Bright Sky Cyan #0ea5e9
COLOR_EMERALD = RGBColor(16, 185, 129)   # Vibrant Emerald #10b981
COLOR_AMBER = RGBColor(245, 158, 11)     # Amber #f59e0b
COLOR_RED = RGBColor(239, 68, 68)        # Crimson #ef4444

COLOR_TEXT_PRIMARY = RGBColor(248, 250, 252) # Crisp White
COLOR_TEXT_MUTED = RGBColor(148, 163, 184)   # Slate Muted
COLOR_TEXT_SUBTLE = RGBColor(203, 213, 225)  # Light Slate

def set_slide_background(slide, color=COLOR_BG):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header(slide, title_text, subtitle_text, tag="WYZMINDZ SOLUTIONS // CYBERSECURITY"):
    # Tag
    t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
    tf_t = t_box.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = tag.upper()
    p_t.font.size = Pt(10)
    p_t.font.bold = True
    p_t.font.color.rgb = COLOR_CYAN
    p_t.font.name = "Arial"

    # Title & Subtitle in one clean block
    h_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.65), Inches(11.7), Inches(0.75))
    tf_h = h_box.text_frame
    tf_h.word_wrap = True
    
    p1 = tf_h.paragraphs[0]
    p1.text = title_text
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT_PRIMARY
    p1.font.name = "Arial"
    
    if subtitle_text:
        p2 = tf_h.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_TEXT_MUTED
        p2.font.name = "Arial"
        p2.space_before = Pt(3)

def add_clean_card(slide, left, top, width, height, title, items, badge_text="", border_color=COLOR_CARD_BORDER, bg_color=COLOR_SURFACE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)
    
    p = tf.paragraphs[0]
    if badge_text:
        p.text = f"[{badge_text}]  {title}"
    else:
        p.text = title
    p.font.size = Pt(13.5)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_PRIMARY
    p.font.name = "Arial"
    
    for item in items:
        p_item = tf.add_paragraph()
        p_item.text = item
        p_item.font.size = Pt(10.5)
        p_item.font.color.rgb = COLOR_TEXT_SUBTLE
        p_item.font.name = "Arial"
        p_item.space_before = Pt(6)
        
    return shape

def create_3_slide_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 1: OBJECTIVE & TRANSFORMATION (BEFORE VS NOW)
    # ══════════════════════════════════════════════════════════════════════════
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1, COLOR_BG)
    add_header(
        s1, 
        "Enterprise VAPT Automation Platform", 
        "Objective: Streamline vulnerability management by transforming manual CLI scanning into an automated assessment suite.",
        "SLIDE 01 // OBJECTIVE & OPERATIONAL TRANSFORMATION"
    )

    # Left Column: BEFORE (The Problem)
    before_items = [
        "❌ Siloed CLI Tools: Running Nmap, ZAP, and SQLmap in isolated terminal tabs with complex manual flags.",
        "❌ Manual Single-Host Scans: No Target Groups; IP addresses had to be typed and scanned one by one.",
        "❌ No Automated Scheduling: Scans were reactive; required writing fragile crontab shell scripts.",
        "❌ 3–5 Hours per Audit Report: Copying raw terminal outputs into Word documents manually.",
        "❌ Data Loss & No Trends: No persistent database to track historical vulnerabilities or open port regressions."
    ]
    add_clean_card(
        s1, Inches(0.8), Inches(1.6), Inches(5.75), Inches(4.3),
        "BEFORE: Fragmented & Manual Workflow",
        before_items,
        badge_text="LEGACY STATE",
        border_color=COLOR_RED,
        bg_color=RGBColor(28, 18, 28)
    )

    # Right Column: NOW (The Solution)
    now_items = [
        "⚡ Unified Web Console: Single-pane dashboard orchestrating 7+ multi-layer security engines.",
        "⚡ Automated Target Groups: Group servers/subnets into asset pools that scan sequentially host-by-host.",
        "⚡ Intelligent APScheduler: Daily (10:00 AM), Weekly, and Monthly recurring jobs with GUI reschedule controls.",
        "⚡ 1-Click Executive Reports: Instant corporate PDF & HTML audits with TOC, OWASP Top 10, and sign-offs.",
        "⚡ Centralized Data Lake: SQLite database indexing all scan runs, port matrices, and severity analytics."
    ]
    add_clean_card(
        s1, Inches(6.75), Inches(1.6), Inches(5.75), Inches(4.3),
        "NOW: Automated Single-Pane Platform",
        now_items,
        badge_text="TRANSFORMED",
        border_color=COLOR_EMERALD,
        bg_color=RGBColor(14, 30, 36)
    )

    # Footer Metadata Bar
    foot1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.85))
    foot1.fill.solid()
    foot1.fill.fore_color.rgb = COLOR_SURFACE
    foot1.line.color.rgb = COLOR_CARD_BORDER
    ftf1 = foot1.text_frame
    ftf1.margin_left = Inches(0.3)
    ftf1.margin_top = Inches(0.15)
    
    fp1 = ftf1.paragraphs[0]
    fp1.text = "Organization: Wyzmindz Solutions   •   Audit Done By: Santhosh M (Network Admin)   •   Approved By: Leo Antony Charles (IT Manager)"
    fp1.font.size = Pt(11)
    fp1.font.bold = True
    fp1.font.color.rgb = COLOR_CYAN

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 2: TOOLS USED & HOW THEY WORK (ARSENAL)
    # ══════════════════════════════════════════════════════════════════════════
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2, COLOR_BG)
    add_header(
        s2, 
        "Security Tools Arsenal & Operational Usage", 
        "Comprehensive multi-layer security coverage across network infrastructure, web applications, and databases.",
        "SLIDE 02 // SECURITY ENGINES & USAGE"
    )

    tools_grid = [
        ("Nmap 7.94", "Network & Port Discovery", ["Identifies live hosts, open TCP/UDP ports, OS details, and service versions across subnets."], COLOR_VIOLET),
        ("OpenVAS (GVM)", "Infrastructure CVE Auditing", ["Performs in-depth vulnerability assessments against OS services, SSL ciphers, and unpatched software."], COLOR_RED),
        ("OWASP ZAP", "Web Application DAST", ["Automates active/passive web security testing, AJAX spidering, and OWASP Top 10 injection audits."], COLOR_CYAN),
        ("Nuclei", "Fast Template & CVE Fuzzing", ["Executes high-speed YAML community templates for rapid zero-day and emerging exploit detection."], COLOR_EMERALD),
        ("Nikto", "Web Server Misconfiguration", ["Scans for 6,700+ dangerous files, outdated server packages, default scripts, and insecure headers."], COLOR_AMBER),
        ("Gobuster & FFuF", "Directory & API Brute-force", ["High-speed wordlist fuzzing to uncover hidden admin panels, backup files, and unlisted API endpoints."], COLOR_CYAN),
        ("SQLmap", "Database Injection Audit", ["Automates detection and validation of SQL injection flaws and database data exfiltration risks."], COLOR_VIOLET),
    ]

    # Render in 2 rows (Top row: 4 tools, Bottom row: 3 tools)
    for i, (name, role, desc_list, color) in enumerate(tools_grid):
        if i < 4:
            x = Inches(0.8 + (i * 2.98))
            y = Inches(1.6)
            w = Inches(2.85)
            h = Inches(2.55)
        else:
            x = Inches(0.8 + ((i - 4) * 3.98))
            y = Inches(4.35)
            w = Inches(3.8)
            h = Inches(2.55)
            
        items = [f"Role: {role}"] + desc_list
        add_clean_card(s2, x, y, w, h, name, items, border_color=color)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 3: AUTOMATED PIPELINE & BUSINESS VALUE (ROI)
    # ══════════════════════════════════════════════════════════════════════════
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3, COLOR_BG)
    add_header(
        s3, 
        "End-to-End Workflow & Business Impact (ROI)", 
        "How automated orchestration delivers measurable operational efficiency and robust compliance readiness.",
        "SLIDE 03 // WORKFLOW & ORGANIZATIONAL VALUE"
    )

    # Top Half: 4-Step Pipeline
    pipe_steps = [
        ("1. Define Target Scope", "Select single host or Target Group (e.g. DMZ, Cloud, Internal LAN)."),
        ("2. Set Timing & Recurrence", "Choose Daily (10:00 AM), Weekly, or Monthly automated APScheduler cron."),
        ("3. Asynchronous Scanning", "FastAPI background workers execute multi-host scans without blocking the UI."),
        ("4. Executive PDF Report", "Instantly generate branded audit reports with TOC, OWASP matrices & sign-offs.")
    ]
    for i, (title, desc) in enumerate(pipe_steps):
        x = Inches(0.8 + (i * 2.98))
        add_clean_card(s3, x, Inches(1.6), Inches(2.85), Inches(2.3), title, [desc], border_color=COLOR_CYAN)

    # Bottom Half: 4 Business ROI Impact Cards
    roi_cards = [
        ("⏱️ 85% Time Savings", "Eliminates repetitive manual scanning and report writing; audit turnaround drops from 4 hours to < 2 minutes.", COLOR_EMERALD),
        ("🛡️ Proactive Security", "Continuous scheduled audits detect misconfigurations and newly exposed ports before attackers can exploit them.", COLOR_CYAN),
        ("💰 Zero Licensing Cost", "100% built on top-tier open-source engines wrapped in an institutionalized automation architecture.", COLOR_VIOLET),
        ("📋 Audit & Compliance Ready", "Produces standardized, executive-level documentation aligned with ISO 27001, SOC 2, and OWASP standards.", COLOR_AMBER)
    ]
    for i, (title, desc, color) in enumerate(roi_cards):
        x = Inches(0.8 + (i * 2.98))
        add_clean_card(s3, x, Inches(4.1), Inches(2.85), Inches(2.8), title, [desc], border_color=color)

    # Save presentation
    output_path = "/root/vapt-dashboard/Wyzmindz_VAPT_Platform_Presentation.pptx"
    prs.save(output_path)
    print(f"Clean 3-Slide Presentation successfully generated: {output_path}")

if __name__ == "__main__":
    create_3_slide_presentation()
